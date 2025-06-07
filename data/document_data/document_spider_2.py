import requests
from bs4 import BeautifulSoup
import os
import json
import time
import hashlib
from urllib.parse import urljoin, urlparse
from concurrent.futures import ThreadPoolExecutor, as_completed
from queue import Queue
import threading
from urllib.parse import unquote
import posixpath
# 爬取不同的学院，选取不同的设置
'''
# 计算机学院
BASE_URL = "https://cc.nankai.edu.cn"
SAVE_DIR = "document_college_data/cc_data"
URL_MAP_FILE = "document_url_map/url_map_cc.json"
MAX_PAGES = 100000
NUM_THREADS = 10
'''

'''
# 金融学院
BASE_URL = "https://finance.nankai.edu.cn/"
SAVE_DIR = "document_college_data/finance_data"
URL_MAP_FILE = "document_url_map/url_map_finance.json"
MAX_PAGES = 100000
NUM_THREADS = 10
'''

'''
# 文学院
BASE_URL = "https://wxy.nankai.edu.cn/"
SAVE_DIR = "document_college_data/wxy_data"
URL_MAP_FILE = "document_url_map/url_map_wxy.json"
MAX_PAGES = 100000
NUM_THREADS = 10
'''

'''
# 哲学院
BASE_URL = "https://phil.nankai.edu.cn/"
SAVE_DIR = "document_college_data/phil_data"
URL_MAP_FILE = "document_url_map/url_map_phil.json"
MAX_PAGES = 100000
NUM_THREADS = 10
'''

'''
# 周恩来政府管理学院
BASE_URL = "https://zfxy.nankai.edu.cn/"
SAVE_DIR = "document_college_data/zfxy_data"
URL_MAP_FILE = "document_url_map/url_map_zfxy.json"
MAX_PAGES = 100000
NUM_THREADS = 10
'''

'''
# 马克思主义学院
BASE_URL = "https://cz.nankai.edu.cn/"
SAVE_DIR = "document_college_data/cz_data"
URL_MAP_FILE = "document_url_map/url_map_cz.json"
MAX_PAGES = 100000
NUM_THREADS = 10
'''

'''
# 商学院
BASE_URL = "https://bs.nankai.edu.cn/"
SAVE_DIR = "document_college_data/bs_data"
URL_MAP_FILE = "document_url_map/url_map_bs.json"
MAX_PAGES = 100000
NUM_THREADS = 10
'''

'''
# 数学科学学院
BASE_URL = "https://math.nankai.edu.cn/"
SAVE_DIR = "document_college_data/math_data"
URL_MAP_FILE = "document_url_map/url_map_math.json"
MAX_PAGES = 100000
NUM_THREADS = 10
'''

'''
# 化学学院
BASE_URL = "https://chem.nankai.edu.cn/"
SAVE_DIR = "document_college_data/chem_data"
URL_MAP_FILE = "document_url_map/url_map_chem.json"
MAX_PAGES = 100000
NUM_THREADS = 10
'''

'''
# 医学院
BASE_URL = "https://medical.nankai.edu.cn/"
SAVE_DIR = "document_college_data/medical_data"
URL_MAP_FILE = "document_url_map/url_map_medical.json"
MAX_PAGES = 100000
NUM_THREADS = 10
'''

'''
# 人工智能学院
BASE_URL = "https://ai.nankai.edu.cn/"
SAVE_DIR = "document_college_data/ai_data"
URL_MAP_FILE = "document_url_map/url_map_ai.json"
MAX_PAGES = 100000
NUM_THREADS = 10
'''

'''
# 电子信息与光学工程学院
BASE_URL = "https://ceo.nankai.edu.cn/"
SAVE_DIR = "document_college_data/ceo_data"
URL_MAP_FILE = "document_url_map/url_map_ceo.json"
MAX_PAGES = 100000
NUM_THREADS = 10
'''

'''
# 汉语言文化学院
BASE_URL = "https://hyxy.nankai.edu.cn/"
SAVE_DIR = "document_college_data/hyxy_data"
URL_MAP_FILE = "document_url_map/url_map_hyxy.json"
MAX_PAGES = 100000
NUM_THREADS = 10
'''

'''
# 药学院
BASE_URL = "https://pharmacy.nankai.edu.cn/"
SAVE_DIR = "document_college_data/pharmacy_data"
URL_MAP_FILE = "document_url_map/url_map_pharmacy.json"
MAX_PAGES = 100000
NUM_THREADS = 10
'''

'''
# 材料科学与工程学院
BASE_URL = "https://mse.nankai.edu.cn/"
SAVE_DIR = "document_college_data/mse_data"
URL_MAP_FILE = "document_url_map/url_map_mse.json"
MAX_PAGES = 100000
NUM_THREADS = 10
'''

'''
# 社会学院
BASE_URL = "https://shxy.nankai.edu.cn/"
SAVE_DIR = "document_college_data/shxy_data"
URL_MAP_FILE = "document_url_map/url_map_shxy.json"
MAX_PAGES = 100000
NUM_THREADS = 10
'''

'''
# 历史学院
BASE_URL = "https://history.nankai.edu.cn/"
SAVE_DIR = "document_college_data/history_data"
URL_MAP_FILE = "document_url_map/url_map_history.json"
MAX_PAGES = 100000
NUM_THREADS = 10
'''

'''
# 法学院
BASE_URL = "https://law.nankai.edu.cn/"
SAVE_DIR = "document_college_data/law_data"
URL_MAP_FILE = "document_url_map/url_map_law.json"
MAX_PAGES = 100000
NUM_THREADS = 10
'''

'''
# 外国语学院
BASE_URL = "https://sfs.nankai.edu.cn/"
SAVE_DIR = "document_college_data/sfs_data"
URL_MAP_FILE = "document_url_map/url_map_sfs.json"
MAX_PAGES = 100000
NUM_THREADS = 10
'''

'''
# 经济学院
BASE_URL = "https://economics.nankai.edu.cn/"
SAVE_DIR = "document_college_data/economics_data"
URL_MAP_FILE = "document_url_map/url_map_economics.json"
MAX_PAGES = 100000
NUM_THREADS = 10
'''

'''
# 统计与数据科学学院
BASE_URL = "https://stat.nankai.edu.cn/"
SAVE_DIR = "document_college_data/stat_data"
URL_MAP_FILE = "document_url_map/url_map_stat.json"
MAX_PAGES = 100000
NUM_THREADS = 10
'''

'''
# 物理科学学院
BASE_URL = "https://physics.nankai.edu.cn/"
SAVE_DIR = "document_college_data/physics_data"
URL_MAP_FILE = "document_url_map/url_map_physics.json"
MAX_PAGES = 100000
NUM_THREADS = 10
'''

'''
# 生命科学学院
BASE_URL = "https://sky.nankai.edu.cn/"
SAVE_DIR = "document_college_data/sky_data"
URL_MAP_FILE = "document_url_map/url_map_sky.json"
MAX_PAGES = 100000
NUM_THREADS = 10
'''

'''
# 网络空间安全学院
BASE_URL = "https://cyber.nankai.edu.cn/"
SAVE_DIR = "document_college_data/cyber_data"
URL_MAP_FILE = "document_url_map/url_map_cyber.json"
MAX_PAGES = 100000
NUM_THREADS = 10
'''

'''
# 软件学院
BASE_URL = "https://cs.nankai.edu.cn/"
SAVE_DIR = "document_college_data/cs_data"
URL_MAP_FILE = "document_url_map/url_map_cs.json"
MAX_PAGES = 100000
NUM_THREADS = 10
'''

'''
# 旅游与服务学院
BASE_URL = "https://tas.nankai.edu.cn/"
SAVE_DIR = "document_college_data/tas_data"
URL_MAP_FILE = "document_url_map/url_map_tas.json"
MAX_PAGES = 100000
NUM_THREADS = 10
'''

'''
# 新闻传播学院
BASE_URL = "https://jc.nankai.edu.cn/"
SAVE_DIR = "document_college_data/jc_data"
URL_MAP_FILE = "document_url_map/url_map_jc.json"
MAX_PAGES = 100000
NUM_THREADS = 10
'''


# 环境科学与工程学院
BASE_URL = "https://env.nankai.edu.cn/"
SAVE_DIR = "document_college_data/env_data"
URL_MAP_FILE = "document_url_map/url_map_env.json"
MAX_PAGES = 100000
NUM_THREADS = 10


# 爬取支持的文档类型
ALLOWED_FILE_EXTENSIONS = {
    '.pdf', '.doc', '.docx', '.xls', '.xlsx', '.ppt', '.pptx'
}
ALLOWED_MIME_TYPES = {
    "application/pdf",
    "application/msword",
    "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
    "application/vnd.ms-excel",
    "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
    "application/vnd.ms-powerpoint",
    "application/vnd.openxmlformats-officedocument.presentationml.presentation"
}

# 初始化目录
os.makedirs(SAVE_DIR, exist_ok=True)
os.makedirs(os.path.dirname(URL_MAP_FILE), exist_ok=True)

# 多线程共享资源
visited = set()
visited_lock = threading.Lock()

url_map = {}
url_map_lock = threading.Lock()

page_count = 0
page_count_lock = threading.Lock()

url_queue = Queue()
url_queue.put(BASE_URL)


def extract_filename_from_url(url):
    path = urlparse(url).path
    return unquote(os.path.basename(path))

# 提取不同文件的文件名
def extract_document_title(file_path):
    try:
        ext = os.path.splitext(file_path)[1].lower()
        # pdf标题提取逻辑
        if ext == '.pdf':
            import pdfplumber
            try:
                with pdfplumber.open(file_path) as pdf:
                    if not pdf.pages:
                        return None
                    page = pdf.pages[0]
                    chars = page.chars
                    if not chars:
                        return None
                    
                    max_size = max(char['size'] for char in chars)
                    max_chars = [c for c in chars if c['size'] == max_size]

                    # 按y坐标从上到下、x坐标从左到右排序，保证顺序合理
                    max_chars.sort(key=lambda c: (c['top'], c['x0']))

                    title = ''.join(c['text'] for c in max_chars).strip()
                    return title[:200] if title else None
            except Exception as e:
                print(f"Failed to extract title from {file_path}: {e}")
                return None
        # docx标题提取逻辑
        elif ext == '.docx':
            from docx import Document
            doc = Document(file_path)
            for para in doc.paragraphs:
                if para.text.strip():
                    return para.text.strip()[:100]

        elif ext in ['.xls', '.xlsx']:
            from openpyxl import load_workbook
            wb = load_workbook(file_path, read_only=True, data_only=True)
            sheet = wb.active
            for row in sheet.iter_rows(values_only=True, max_row=5):
                for cell in row:
                    if cell and isinstance(cell, str):
                        return cell.strip()[:100]

        elif ext == '.pptx':
            from pptx import Presentation
            ppt = Presentation(file_path)
            for slide in ppt.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        return shape.text.strip()[:100]

    except Exception as e:
        print(f"Failed to extract title from {file_path}: {e}")
    return None

# 判断是否是有效的文档
def is_valid_document(url, content_type):
    parsed = urlparse(url)
    ext = os.path.splitext(parsed.path)[1].lower()
    if ext in ALLOWED_FILE_EXTENSIONS:
        return True
    if any(mt in content_type for mt in ALLOWED_MIME_TYPES):
        return True
    return False

def fetch_url(url):
    try:
        response = requests.get(url, timeout=10, headers={'User-Agent': 'Mozilla/5.0'})
        response.raise_for_status()
        return response
    except requests.RequestException as e:
        print(f"Failed to fetch {url}: {e}")
        return None

def parse_links(html, base_url):
    soup = BeautifulSoup(html, "html.parser")
    links = set()
    for tag in soup.find_all("a", href=True):
        href = tag['href'].replace('\\', '/')
        full_url = urljoin(base_url, href)
        parsed_url = urlparse(full_url)
        if parsed_url.netloc == urlparse(base_url).netloc:
            links.add(full_url.split("#")[0])
    return links

def get_doc_filename(url):
    url_hash = hashlib.md5(url.encode('utf-8')).hexdigest()
    ext = os.path.splitext(urlparse(url).path)[1].lower()
    return os.path.join(SAVE_DIR, f"{url_hash}{ext}")

def save_file(content, file_path):
    with open(file_path, "wb") as f:
        f.write(content)

# 核心并行函数
def worker():
    global page_count
    while True:
        if page_count >= MAX_PAGES:
            break

        try:
            url = url_queue.get(timeout=1)
        except:
            break

        with visited_lock:
            if url in visited:
                continue
            visited.add(url)

        print(f"[Thread-{threading.get_ident()}] Checking: {url}")
        response = fetch_url(url)
        if not response:
            continue

        content_type = response.headers.get("Content-Type", "").split(";")[0].strip().lower()
        if is_valid_document(url, content_type):
            file_path = get_doc_filename(url)
            save_file(response.content, file_path)

            title = extract_document_title(file_path)

            with url_map_lock:
                url_map[file_path] = {
                    "url": url,
                    "original_filename": extract_filename_from_url(url),
                    "title": title if title else ""
                }

            with page_count_lock:
                page_count += 1
                if page_count >= MAX_PAGES:
                    break

        # 如果是html，就继续提取链接
        elif "text/html" in content_type:
            response.encoding = response.apparent_encoding
            links = parse_links(response.text, BASE_URL)
            with visited_lock:
                for link in links:
                    if link not in visited:
                        url_queue.put(link)

        time.sleep(0.01)

# 并行爬取主函数
def scrape_parallel():
    with ThreadPoolExecutor(max_workers=NUM_THREADS) as executor:
        futures = [executor.submit(worker) for _ in range(NUM_THREADS)]
        for f in as_completed(futures):
            pass

    with open(URL_MAP_FILE, "w", encoding="utf-8") as f:
        json.dump(url_map, f, indent=2, ensure_ascii=False)

    print(f"Downloaded {page_count} documents.")
    print(f"URL mapping saved to {URL_MAP_FILE}")

if __name__ == "__main__":
    scrape_parallel()